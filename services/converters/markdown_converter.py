from typing import Dict, Any, List, Optional, BinaryIO, Tuple, Union
from pathlib import Path
import mammoth
import PyPDF2
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.oxml.ns import qn
from docx.shared import Inches
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from docx.enum.text import WD_ALIGN_PARAGRAPH
from lxml.etree import _Element  # This is the correct import
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from openpyxl import load_workbook
import tempfile
import os
import io
import base64
from loguru import logger
import asyncio
import markdown
import re
import json
import unicodedata
from dataclasses import dataclass
from datetime import datetime
from core.exceptions import FileConversionException
from models.file_conversion_models import FileType, FileMetadata, ConversionWarning

@dataclass
class DocumentElement:
    """Represents a semantic element in the document"""
    type: str  # heading, paragraph, list, table, image, etc.
    content: Union[str, List[str], Dict]
    metadata: Dict[str, Any]
    level: Optional[int] = None
    children: List['DocumentElement'] = None

    def __post_init__(self):
        if self.children is None:
            self.children = []

class DocumentStructure:
    """Manages document structure and hierarchy"""
    def __init__(self):
        self.elements: List[DocumentElement] = []
        self.current_section: Optional[DocumentElement] = None
        self.section_stack: List[DocumentElement] = []
    
    def add_element(self, element: DocumentElement):
        """Add element maintaining hierarchy"""
        if element.type == 'heading':
            self._handle_heading(element)
        elif self.current_section:
            self.current_section.children.append(element)
        else:
            self.elements.append(element)
    
    def _handle_heading(self, heading: DocumentElement):
        """Handle heading hierarchy"""
        while (self.section_stack and 
               self.section_stack[-1].level is not None and 
               self.section_stack[-1].level >= heading.level):
            self.section_stack.pop()
        
        if self.section_stack:
            self.section_stack[-1].children.append(heading)
        else:
            self.elements.append(heading)
        
        self.section_stack.append(heading)
        self.current_section = heading

class BaseConverter:
    """Base class for all file converters"""
    
    async def convert(self, file_content: bytes) -> Tuple[str, FileMetadata]:
        """Convert file content to markdown"""
        raise NotImplementedError
    
    def _create_temp_file(self, content: bytes, suffix: str) -> str:
        """Create a temporary file with the given content and suffix"""
        temp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        try:
            temp.write(content)
            temp.flush()
            return temp.name
        finally:
            temp.close()

    def _clean_text(self, text: str) -> str:
        """Clean and normalize text content"""
        # Remove control characters except newlines
        text = ''.join(char if char == '\n' or not unicodedata.category(char).startswith('C') else ' ' for char in text)
        # Normalize whitespace
        text = re.sub(r'[^\S\n]+', ' ', text)
        # Normalize newlines
        text = re.sub(r'\n{3,}', '\n\n', text)
        # Remove whitespace at start/end of lines
        text = '\n'.join(line.strip() for line in text.split('\n'))
        return text.strip()
    
    def _encode_image(self, image_data: bytes, image_type: str = 'png') -> str:
        """Encode image as base64 for markdown"""
        try:
            encoded = base64.b64encode(image_data).decode('utf-8')
            return f"data:image/{image_type};base64,{encoded}"
        except Exception as e:
            logger.error(f"Image encoding error: {str(e)}")
            return ""

    def _structure_to_markdown(self, structure: DocumentStructure) -> str:
        """Convert document structure to markdown"""
        def process_element(element: DocumentElement, level: int = 0) -> List[str]:
            lines = []
            
            if element.type == 'heading':
                lines.append(f"{'#' * element.level} {element.content}")
            elif element.type == 'paragraph':
                lines.append(str(element.content))
            elif element.type == 'list':
                prefix = '  ' * level
                if isinstance(element.content, list):
                    for item in element.content:
                        lines.append(f"{prefix}* {item}")
            elif element.type == 'table':
                lines.extend(self._format_table(element.content))
            elif element.type == 'image':
                if isinstance(element.content, dict) and 'data' in element.content:
                    image_data = self._encode_image(element.content['data'])
                    alt_text = element.metadata.get('alt_text', 'Image')
                    lines.append(f"![{alt_text}]({image_data})")
            elif element.type == 'code':
                lang = element.metadata.get('language', '')
                lines.append(f"```{lang}")
                lines.append(element.content)
                lines.append("```")
            
            # Process children
            if element.children:
                child_lines = []
                for child in element.children:
                    child_lines.extend(process_element(child, level + 1))
                lines.extend(child_lines)
            
            return lines
        
        # Process all elements
        markdown_lines = []
        for element in structure.elements:
            markdown_lines.extend(process_element(element))
        
        return '\n\n'.join(markdown_lines)
    
    def _format_table(self, table_data: List[List[str]]) -> List[str]:
        """Format table data as markdown"""
        if not table_data or not table_data[0]:
            return []
            
        # Calculate column widths
        col_widths = [0] * len(table_data[0])
        for row in table_data:
            for i, cell in enumerate(row):
                col_widths[i] = max(col_widths[i], len(str(cell)))
        
        # Create header
        header = "| " + " | ".join(f"{str(cell):<{width}}" for cell, width in zip(table_data[0], col_widths)) + " |"
        separator = "|" + "|".join(f":{'-'*(width-2)}:" for width in col_widths) + "|"
        
        # Create rows
        rows = []
        for row in table_data[1:]:
            formatted_row = "| " + " | ".join(f"{str(cell):<{width}}" for cell, width in zip(row, col_widths)) + " |"
            rows.append(formatted_row)
        
        return [header, separator] + rows
    
class PDFConverter(BaseConverter):
    """Enhanced PDF converter with semantic structure preservation"""
    
    def _extract_logical_sections(self, text: str) -> List[DocumentElement]:
        """Extract logical sections from text using semantic analysis"""
        sections = []
        lines = text.split('\n')
        current_block = []
        current_type = 'paragraph'
        
        def flush_block():
            if not current_block:
                return None
            content = ' '.join(current_block)
            # Analyze content type
            if re.match(r'^[A-Z][A-Z\s]{8,}$', content):  # All caps = heading
                return DocumentElement('heading', content, {'style': 'uppercase'}, level=2)
            elif re.match(r'^\d+[\.\)]\s', content):  # Numbered list
                return DocumentElement('list', content, {'list_type': 'ordered'})
            return DocumentElement('paragraph', content, {})
        
        for line in lines:
            line = line.strip()
            if not line:
                if element := flush_block():
                    sections.append(element)
                current_block = []
                continue
                
            # Detect headings by font size/style (if available in PDF metadata)
            # Detect lists
            if re.match(r'^\s*[\-\*\d]+[\.\)]\s', line):
                if current_block and current_type != 'list':
                    if element := flush_block():
                        sections.append(element)
                    current_block = []
                current_type = 'list'
            
            current_block.append(line)
        
        if element := flush_block():
            sections.append(element)
            
        return sections

    def _extract_images(self, page: Any) -> List[Dict]:
        """Extract images from PDF page with metadata"""
        images = []
        try:
            if '/XObject' in page:
                xObject = page['/XObject'].get_object()
                for obj in xObject:
                    if xObject[obj]['/Subtype'] == '/Image':
                        try:
                            image_data = self._extract_image_data(xObject[obj])
                            if image_data:
                                image_info = {
                                    'data': image_data,
                                    'metadata': {
                                        'width': xObject[obj].get('/Width', 0),
                                        'height': xObject[obj].get('/Height', 0),
                                        'format': xObject[obj].get('/Filter', 'Unknown'),
                                        'bits': xObject[obj].get('/BitsPerComponent', 0)
                                    }
                                }
                                images.append(image_info)
                        except Exception as e:
                            logger.warning(f"Failed to extract image: {str(e)}")
        except Exception as e:
            logger.error(f"Error processing images: {str(e)}")
        return images

    def _extract_image_data(self, image_object: Any) -> Optional[bytes]:
        """Extract raw image data from PDF image object"""
        try:
            if image_object['/Filter'] == '/FlateDecode':
                data = image_object.get_data()  # Changed from getData to get_data
                return data
            elif image_object['/Filter'] == '/DCTDecode':
                return image_object._data
            elif image_object['/Filter'] == '/JPXDecode':
                return image_object._data
        except Exception as e:
            logger.error(f"Image data extraction error: {str(e)}")
        return None

    def _extract_text_with_formatting(self, page: Any) -> Tuple[str, Dict[str, Any]]:
        """Extract text while preserving formatting information"""
        text = page.extract_text()
        formatting = {}
        try:
            # Try to extract font information
            if '/Resources' in page and '/Font' in page['/Resources']:
                fonts = page['/Resources']['/Font']
                formatting['fonts'] = [str(font) for font in fonts]
            
            # Try to extract text positioning
            if '/Contents' in page:
                # Changed from getObject() to get_object()
                content = page['/Contents'].get_object()
                if hasattr(content, 'get_data'):  # Changed from getData to get_data
                    ops = content.get_data()  # Changed from getData to get_data
                    # Parse PDF operators for text positioning
                    formatting['positioning'] = self._parse_text_positioning(ops)
        except Exception as e:
            logger.warning(f"Error extracting formatting: {str(e)}")
        
        return text, formatting

    def _detect_tables(self, page: Any) -> List[Dict[str, Any]]:
        """Detect and extract tables from PDF"""
        tables = []
        try:
            # Use positioning and spacing to detect table structures
            text, formatting = self._extract_text_with_formatting(page)
            # Implement table detection logic based on text positioning
            # This is a simplified version - real implementation would be more complex
            if 'positioning' in formatting:
                potential_tables = self._find_table_structures(text, formatting['positioning'])
                tables.extend(potential_tables)
        except Exception as e:
            logger.warning(f"Table detection error: {str(e)}")
        return tables

    async def convert(self, file_content: bytes) -> Tuple[str, FileMetadata]:
        temp_path = self._create_temp_file(file_content, '.pdf')
        try:
            reader = PyPDF2.PdfReader(temp_path)
            doc_structure = DocumentStructure()
            images_found = []
            tables_found = []
            
            # Extract document metadata
            if reader.metadata:
                metadata_element = DocumentElement(
                    'metadata',
                    reader.metadata,
                    {'source': 'pdf_metadata'}
                )
                doc_structure.add_element(metadata_element)
            
            # Process each page
            for page_num, page in enumerate(reader.pages, 1):
                # Add page marker
                page_element = DocumentElement(
                    'heading',
                    f"Page {page_num}",
                    {'type': 'page_marker'},
                    level=1
                )
                doc_structure.add_element(page_element)
                
                # Extract text with formatting
                text, formatting = self._extract_text_with_formatting(page)
                
                # Extract and process images
                page_images = self._extract_images(page)
                for img_num, img in enumerate(page_images, 1):
                    image_element = DocumentElement(
                        'image',
                        img['data'],
                        {
                            'page': page_num,
                            'image_num': img_num,
                            **img['metadata']
                        }
                    )
                    doc_structure.add_element(image_element)
                    images_found.append(img)
                
                # Detect and extract tables
                page_tables = self._detect_tables(page)
                for table_num, table in enumerate(page_tables, 1):
                    table_element = DocumentElement(
                        'table',
                        table['data'],
                        {
                            'page': page_num,
                            'table_num': table_num,
                            'headers': table.get('headers', [])
                        }
                    )
                    doc_structure.add_element(table_element)
                    tables_found.append(table)
                
                # Process text content
                if text.strip():
                    sections = self._extract_logical_sections(text)
                    for section in sections:
                        doc_structure.add_element(section)
            
            # Convert structure to markdown
            markdown_content = self._structure_to_markdown(doc_structure)
            
            metadata = FileMetadata(
                filename=os.path.basename(temp_path),
                size_bytes=os.path.getsize(temp_path),
                file_type=FileType.PDF,
                pages=len(reader.pages),
                images_count=len(images_found),
                tables_count=len(tables_found),
                equations_count=None,
                semantic_structure=doc_structure.elements
            )
            
            return markdown_content, metadata
            
        except Exception as e:
            logger.error(f"PDF conversion error: {str(e)}")
            raise FileConversionException(f"Failed to convert PDF: {str(e)}")
        finally:
            try:
                os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Error cleaning up temporary PDF file: {str(e)}")

class DocxConverter(BaseConverter):
    """Enhanced DOCX converter with style and structure preservation"""
    
    def _parse_paragraph_style(self, paragraph) -> Dict[str, Any]:
        """Extract paragraph style information"""
        style_info = {
            'name': paragraph.style.name if paragraph.style else 'Normal',
            'alignment': str(paragraph.alignment) if hasattr(paragraph, 'alignment') else None,
            'indentation': {},
            'spacing': {},
        }
        
        try:
            if paragraph._element.pPr is not None:
                if paragraph._element.pPr.ind is not None:
                    ind = paragraph._element.pPr.ind
                    style_info['indentation'].update({
                        'left': ind.left.__str__() if ind.left else None,
                        'right': ind.right.__str__() if ind.right else None,
                        'firstLine': ind.firstLine.__str__() if ind.firstLine else None,
                    })
                
                if paragraph._element.pPr.spacing is not None:
                    spacing = paragraph._element.pPr.spacing
                    style_info['spacing'].update({
                        'before': spacing.before.__str__() if spacing.before else None,
                        'after': spacing.after.__str__() if spacing.before else None,
                        'line': spacing.line.__str__() if spacing.line else None,
                    })
        except Exception as e:
            logger.warning(f"Error extracting paragraph style: {str(e)}")
        
        return style_info

    def _process_table(self, table) -> DocumentElement:
        """Process DOCX table with formatting"""
        table_data = []
        cell_styles = []
        
        for row in table.rows:
            row_data = []
            row_styles = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                
                # Extract cell style
                style = {
                    'background_color': cell._tc.get_or_add_tcPr().shd_val,
                    'vertical_align': cell._tc.get_or_add_tcPr().vAlign_val,
                    'width': cell.width if hasattr(cell, 'width') else None,
                }
                
                row_data.append(cell_text)
                row_styles.append(style)
            
            table_data.append(row_data)
            cell_styles.append(row_styles)
        
        return DocumentElement(
            'table',
            table_data,
            {
                'styles': cell_styles,
                'has_header': True  # Could be detected based on styling
            }
        )

    def _process_image(self, shape) -> Optional[DocumentElement]:
        """Process embedded image"""
        try:
            image_data = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            image_path = shape._parent._part.related_parts[image_data].image.path
            
            with open(image_path, 'rb') as img_file:
                image_data = img_file.read()
                
            return DocumentElement(
                'image',
                {'data': image_data},
                {
                    'width': shape.width,
                    'height': shape.height,
                    'filename': os.path.basename(image_path)
                }
            )
        except Exception as e:
            logger.warning(f"Error processing image: {str(e)}")
            return None

    async def convert(self, file_content: bytes) -> Tuple[str, FileMetadata]:
        temp_path = self._create_temp_file(file_content, '.docx')
        try:
            doc = Document(temp_path)
            doc_structure = DocumentStructure()
            images_found = []
            tables_found = []
            
            # Document Properties
            if doc.core_properties:
                props = doc.core_properties
                doc_structure.add_element(DocumentElement(
                    'metadata',
                    {
                        'title': props.title,
                        'author': props.author,
                        'created': props.created.isoformat() if props.created else None,
                        'modified': props.modified.isoformat() if props.modified else None,
                        'revision': props.revision
                    },
                    {'source': 'document_properties'}
                ))
            
            # Process content
            for element in doc.element.body:
                if isinstance(element, _Element):  # Document element
                    if isinstance(element, CT_P):  # Paragraph
                        paragraph = Paragraph(element, doc)
                        style = self._parse_paragraph_style(paragraph)
                        
                        if style['name'].startswith('Heading'):
                            level = int(style['name'][-1]) if style['name'][-1].isdigit() else 1
                            doc_structure.add_element(DocumentElement(
                                'heading',
                                paragraph.text,
                                {'style': style},
                                level=level
                            ))
                        else:
                            doc_structure.add_element(DocumentElement(
                                'paragraph',
                                paragraph.text,
                                {'style': style}
                            ))
                            
                    elif isinstance(element, CT_Tbl):  # Table
                        table = Table(element, doc)
                        table_element = self._process_table(table)
                        doc_structure.add_element(table_element)
                        tables_found.append(table_element)
            
            # Process images
            for shape in doc.inline_shapes:
                if shape.type == MSO_SHAPE_TYPE.PICTURE:
                    image_element = self._process_image(shape)
                    if image_element:
                        doc_structure.add_element(image_element)
                        images_found.append(image_element)
            
            # Convert to markdown
            markdown_content = self._structure_to_markdown(doc_structure)
            
            metadata = FileMetadata(
                filename=os.path.basename(temp_path),
                size_bytes=os.path.getsize(temp_path),
                file_type=FileType.DOCX,
                pages=len(doc.sections),
                images_count=len(images_found),
                tables_count=len(tables_found),
                equations_count=sum(1 for p in doc.paragraphs if 'math' in p._element.xml),
                semantic_structure=doc_structure.elements
            )
            
            return markdown_content, metadata
            
        except Exception as e:
            logger.error(f"DOCX conversion error: {str(e)}")
            raise FileConversionException(f"Failed to convert DOCX: {str(e)}")
        finally:
            try:
                os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Error cleaning up temporary DOCX file: {str(e)}")

class XlsxConverter(BaseConverter):
    """Enhanced XLSX converter with formula and formatting preservation"""
    
    def _detect_table_range(self, sheet) -> Tuple[int, int, int, int]:
        """Detect actual data range in sheet"""
        min_row = min_col = float('inf')
        max_row = max_col = 0
        
        for row in range(1, sheet.max_row + 1):
            for col in range(1, sheet.max_column + 1):
                cell = sheet.cell(row=row, column=col)
                if cell.value is not None:
                    min_row = min(min_row, row)
                    min_col = min(min_col, col)
                    max_row = max(max_row, row)
                    max_col = max(max_col, col)
        
        return (min_row, min_col, max_row, max_col)

    def _process_cell(self, cell) -> Dict[str, Any]:
        """Process cell with value and formatting"""
        cell_info = {
            'value': cell.value,
            'formula': cell.formula if cell.formula else None,
            'format': cell.number_format if cell.number_format != 'General' else None,
            'style': {}
        }
        
        if cell.font:
            cell_info['style']['font'] = {
                'bold': cell.font.bold,
                'italic': cell.font.italic,
                'size': cell.font.size,
                'color': cell.font.color.rgb if cell.font.color else None
            }
        
        if cell.fill:
            cell_info['style']['fill'] = {
                'type': cell.fill.fill_type,
                'color': cell.fill.start_color.rgb if cell.fill.start_color else None
            }
        
        return cell_info

    def _get_merged_cell_ranges(self, sheet) -> List[str]:
        """Get list of merged cell ranges"""
        return [str(merged_cell) for merged_cell in sheet.merged_cells.ranges]

    async def convert(self, file_content: bytes) -> Tuple[str, FileMetadata]:
        temp_path = self._create_temp_file(file_content, '.xlsx')
        try:
            wb = load_workbook(temp_path, data_only=True)
            doc_structure = DocumentStructure()
            tables_found = []
            
            # Workbook properties
            doc_structure.add_element(DocumentElement(
                'metadata',
                {
                    'sheets': wb.sheetnames,
                    'properties': {
                        'title': wb.properties.title,
                        'subject': wb.properties.subject,
                        'creator': wb.properties.creator,
                        'modified': wb.properties.modified.isoformat() if wb.properties.modified else None
                    }
                },
                {'source': 'workbook_properties'}
            ))
            
            for sheet in wb:
                # Add sheet heading
                doc_structure.add_element(DocumentElement(
                    'heading',
                    f"Sheet: {sheet.title}",
                    {'sheet_name': sheet.title},
                    level=1
                ))
                
                # Detect actual data range
                min_row, min_col, max_row, max_col = self._detect_table_range(sheet)
                if min_row == float('inf'):
                    continue
                
                # Get merged cells
                merged_ranges = self._get_merged_cell_ranges(sheet)
                
                # Process data
                table_data = []
                table_styles = []
                
                for row in range(min_row, max_row + 1):
                    row_data = []
                    row_styles = []
                    for col in range(min_col, max_col + 1):
                        cell = sheet.cell(row=row, column=col)
                        cell_info = self._process_cell(cell)
                        row_data.append(cell_info['value'])
                        row_styles.append(cell_info['style'])
                    
                    table_data.append(row_data)
                    table_styles.append(row_styles)
                
                table_element = DocumentElement(
                    'table',
                    table_data,
                    {
                        'sheet': sheet.title,
                        'styles': table_styles,
                        'merged_ranges': merged_ranges,
                        'has_header': True  # Could be detected based on styling
                    }
                )
                
                doc_structure.add_element(table_element)
                tables_found.append(table_element)
            
            # Convert to markdown
            markdown_content = self._structure_to_markdown(doc_structure)
            
            metadata = FileMetadata(
                filename=os.path.basename(temp_path),
                size_bytes=os.path.getsize(temp_path),
                file_type=FileType.XLSX,
                pages=len(wb.sheetnames),
                images_count=0,  # Excel shapes not handled yet
                tables_count=len(tables_found),
                equations_count=None,
                semantic_structure=doc_structure.elements
            )
            
            return markdown_content, metadata
            
        except Exception as e:
            logger.error(f"XLSX conversion error: {str(e)}")
            raise FileConversionException(f"Failed to convert XLSX: {str(e)}")
        finally:
            try:
                os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Error cleaning up temporary XLSX file: {str(e)}")

class PptxConverter(BaseConverter):
    """Enhanced PPTX converter with layout and formatting preservation"""
    
    def _process_shape_text(self, shape) -> Optional[DocumentElement]:
        """Process shape text with formatting"""
        if not hasattr(shape, "text") or not shape.text.strip():
            return None
            
        text_frame = shape.text_frame
        paragraphs = []
        
        for paragraph in text_frame.paragraphs:
            para_info = {
                'text': paragraph.text,
                'level': paragraph.level,
                'alignment': paragraph.alignment,
                'runs': []
            }
            
            for run in paragraph.runs:
                run_info = {
                    'text': run.text,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'size': run.font.size,
                    'color': run.font.color.rgb if run.font.color else None
                }
                para_info['runs'].append(run_info)
            
            paragraphs.append(para_info)
        
        return DocumentElement(
            'text',
            paragraphs,
            {
                'shape_type': shape.shape_type,
                'width': shape.width,
                'height': shape.height
            }
        )

    def _process_table(self, shape) -> Optional[DocumentElement]:
        """Process table with formatting"""
        if not shape.has_table:
            return None
            
        table = shape.table
        table_data = []
        cell_styles = []
        
        for row in table.rows:
            row_data = []
            row_styles = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                
                style = {
                    'fill': cell.fill.type if cell.fill else None,
                    'margin': cell.margin,
                    'vertical_anchor': cell.vertical_anchor
                }
                
                row_data.append(cell_text)
                row_styles.append(style)
            
            table_data.append(row_data)
            cell_styles.append(row_styles)
        
        return DocumentElement(
            'table',
            table_data,
            {
                'styles': cell_styles,
                'has_header': True
            }
        )

    def _process_shape_image(self, shape) -> Optional[DocumentElement]:
        """Process shape image with metadata"""
        try:
            if not hasattr(shape, 'image'):
                return None
                
            image_data = shape.image.blob
            image_info = {
                'data': image_data,
                'filename': shape.image.filename,
                'content_type': shape.image.content_type,
                'size': len(image_data)
            }
            
            return DocumentElement(
                'image',
                image_info,
                {
                    'width': shape.width,
                    'height': shape.height,
                    'position_x': shape.left,
                    'position_y': shape.top
                }
            )
        except Exception as e:
            logger.warning(f"Error processing image shape: {str(e)}")
            return None

    def _process_notes(self, notes_slide) -> Optional[DocumentElement]:
        """Process slide notes"""
        if not notes_slide or not notes_slide.notes_text_frame.text.strip():
            return None
            
        return DocumentElement(
            'notes',
            notes_slide.notes_text_frame.text.strip(),
            {'type': 'speaker_notes'}
        )

    def _get_slide_layout_info(self, slide) -> Dict[str, Any]:
        """Extract slide layout information"""
        layout_info = {
            'type': slide.slide_layout.name,
            'shapes': [],
            'placeholders': []
        }
        
        try:
            for shape in slide.shapes:
                shape_info = {
                    'type': shape.shape_type,
                    'name': shape.name,
                    'width': shape.width,
                    'height': shape.height,
                    'position': (shape.left, shape.top)
                }
                if hasattr(shape, 'placeholder_format'):
                    shape_info['placeholder_type'] = shape.placeholder_format.type
                    layout_info['placeholders'].append(shape_info)
                else:
                    layout_info['shapes'].append(shape_info)
        except Exception as e:
            logger.warning(f"Error extracting layout info: {str(e)}")
        
        return layout_info

    async def convert(self, file_content: bytes) -> Tuple[str, FileMetadata]:
        temp_path = self._create_temp_file(file_content, '.pptx')
        try:
            prs = Presentation(temp_path)
            doc_structure = DocumentStructure()
            images_found = []
            tables_found = []
            
            # Presentation properties
            if prs.core_properties:
                doc_structure.add_element(DocumentElement(
                    'metadata',
                    {
                        'title': prs.core_properties.title,
                        'author': prs.core_properties.author,
                        'created': prs.core_properties.created.isoformat() if prs.core_properties.created else None,
                        'modified': prs.core_properties.modified.isoformat() if prs.core_properties.modified else None,
                        'revision': prs.core_properties.revision,
                        'slides_count': len(prs.slides)
                    },
                    {'source': 'presentation_properties'}
                ))
            
            # Process slides
            for slide_num, slide in enumerate(prs.slides, 1):
                # Add slide marker
                slide_element = DocumentElement(
                    'heading',
                    f"Slide {slide_num}",
                    {
                        'slide_number': slide_num,
                        'layout': self._get_slide_layout_info(slide)
                    },
                    level=1
                )
                doc_structure.add_element(slide_element)
                
                # Process slide content
                for shape in slide.shapes:
                    # Process text
                    if hasattr(shape, "text"):
                        text_element = self._process_shape_text(shape)
                        if text_element:
                            doc_structure.add_element(text_element)
                    
                    # Process tables
                    if shape.has_table:
                        table_element = self._process_table(shape)
                        if table_element:
                            doc_structure.add_element(table_element)
                            tables_found.append(table_element)
                    
                    # Process images
                    image_element = self._process_shape_image(shape)
                    if image_element:
                        doc_structure.add_element(image_element)
                        images_found.append(image_element)
                
                # Process notes
                if slide.has_notes_slide:
                    notes_element = self._process_notes(slide.notes_slide)
                    if notes_element:
                        doc_structure.add_element(notes_element)
                
                # Add slide separator
                doc_structure.add_element(DocumentElement(
                    'separator',
                    '---',
                    {'type': 'slide_break'}
                ))
            
            # Convert to markdown
            markdown_content = self._structure_to_markdown(doc_structure)
            
            metadata = FileMetadata(
                filename=os.path.basename(temp_path),
                size_bytes=os.path.getsize(temp_path),
                file_type=FileType.PPTX,
                pages=len(prs.slides),
                images_count=len(images_found),
                tables_count=len(tables_found),
                equations_count=None,
                semantic_structure=doc_structure.elements
            )
            
            return markdown_content, metadata
            
        except Exception as e:
            logger.error(f"PPTX conversion error: {str(e)}")
            raise FileConversionException(f"Failed to convert PPTX: {str(e)}")
        finally:
            try:
                os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Error cleaning up temporary PPTX file: {str(e)}")


